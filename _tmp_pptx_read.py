# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

path = r"C:/Users/Mr.S/xwechat_files/wxid_r0ty34wac6bt22_7a6c/msg/file/2026-04/Viktor-分享pptx.pptx"
p = Presentation(path)

print(f"Total slides: {len(p.slides)}")
print("=" * 60)

for i, slide in enumerate(p.slides, 1):
    print(f"\n### Slide {i}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs)
                if text.strip():
                    print(text)
        elif shape.shape_type == 19:  # TABLE
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    print(" | ".join(cells))
    # Check notes
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"[NOTES]: {notes}")
