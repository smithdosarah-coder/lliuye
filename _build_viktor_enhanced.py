# -*- coding: utf-8 -*-
"""Insert 2 new slides before the final 'thank you' slide (slide 43)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

PATH = r"D:/claude code/credit_report_agent_work/Viktor-分享_enhanced.pptx"

# Palette
BG        = RGBColor(0x0A, 0x0E, 0x1A)
CARD_BG   = RGBColor(0x14, 0x1B, 0x2D)
ACCENT    = RGBColor(0x6C, 0x8A, 0xFF)
ACCENT2   = RGBColor(0xFF, 0xB6, 0x4E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SUB       = RGBColor(0xAE, 0xB3, 0xC2)

FONT = "Microsoft YaHei"
EMOJI_FONT = "Segoe UI Emoji"


def set_slide_bg(slide, color):
    """Paint solid background on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font_size=14, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_card(slide, left, top, width, height, *, fill=CARD_BG, line=ACCENT,
             line_width_pt=1.0, corner=0.06):
    """Rounded card with subtle accent border."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.adjustments[0] = corner  # subtle rounding
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(line_width_pt)
    # remove shadow
    sppr = shp.fill._xPr  # spPr element
    # Explicitly add effectLst empty to disable inherited shadow
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    # remove existing effectLst
    for el in sppr.findall(qn('a:effectLst')):
        sppr.remove(el)
    effect = etree.SubElement(sppr, qn('a:effectLst'))
    shp.text = ""  # clear
    shp.text_frame.margin_left = Inches(0.25)
    shp.text_frame.margin_right = Inches(0.25)
    shp.text_frame.margin_top = Inches(0.2)
    shp.text_frame.margin_bottom = Inches(0.2)
    shp.text_frame.word_wrap = True
    return shp


def fill_card_title_body(card, title, body, *, title_size=20, body_size=13,
                         body_color=SUB):
    tf = card.text_frame
    tf.clear()
    # title paragraph
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run_t = p1.add_run()
    run_t.text = title
    run_t.font.name = FONT
    run_t.font.size = Pt(title_size)
    run_t.font.bold = True
    run_t.font.color.rgb = WHITE
    # body paragraph
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(10)
    run_b = p2.add_run()
    run_b.text = body
    run_b.font.name = FONT
    run_b.font.size = Pt(body_size)
    run_b.font.color.rgb = body_color


def fill_card_compare(card, title, rows):
    """rows = [('旧', text, color), ('新', text, color), ('位移点', text, color)]"""
    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r = p1.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(19)
    r.font.bold = True
    r.font.color.rgb = WHITE

    for (label, text, color) in rows:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(8)
        r_lbl = p.add_run()
        r_lbl.text = label + "："
        r_lbl.font.name = FONT
        r_lbl.font.size = Pt(12)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = ACCENT
        r_body = p.add_run()
        r_body.text = text
        r_body.font.name = FONT
        r_body.font.size = Pt(12)
        r_body.font.bold = (color == ACCENT2)
        r_body.font.color.rgb = color


def build_slide_A(slide, sw, sh):
    set_slide_bg(slide, BG)

    # eyebrow (top-left)
    add_textbox(slide, Inches(0.5), Inches(0.45),
                Inches(4.0), Inches(0.35),
                "06 / 协作机制",
                font_size=14, bold=True, color=ACCENT)

    # main title
    add_textbox(slide, Inches(0.5), Inches(0.85),
                Inches(9.0), Inches(0.7),
                "Viktor ≠ AI 工具 · 而是团队的 AI 同事",
                font_size=34, bold=True, color=WHITE)

    # subtitle (english)
    add_textbox(slide, Inches(0.5), Inches(1.55),
                Inches(9.0), Inches(0.35),
                "The collaboration mechanics that set Viktor apart",
                font_size=13, italic=True, color=SUB)

    # 2x2 grid of cards
    margin_l = Inches(0.5)
    margin_r = Inches(0.5)
    top_grid = Inches(2.15)
    gap = Inches(0.3)
    grid_w = sw - margin_l - margin_r
    card_w = (grid_w - gap) / 2
    card_h = Inches(1.35)

    cards_data = [
        ("👀  频道可见",
         "执行过程实时在 Slack 频道里可见。team 可以围观、评论、甚至接手。"),
        ("📚  知识沉淀",
         "所有任务对话自动变成团队知识库。可搜索、可回溯——不再“知识留在个人电脑”。"),
        ("✋  随时接管",
         "任意成员 @ 它即中断正在跑的任务。像对真实同事下新指令，不用等任务跑完。"),
        ("🔔  主动协作",
         "不只被动响应，会主动 ping：“周一早上了，上周任务进展如何？”"),
    ]
    for idx, (title, body) in enumerate(cards_data):
        row = idx // 2
        col = idx % 2
        left = margin_l + col * (card_w + gap)
        top = top_grid + row * (card_h + gap)
        c = add_card(slide, left, top, card_w, card_h)
        fill_card_title_body(c, title, body, title_size=19, body_size=13)

    # bottom tagline
    add_textbox(slide, Inches(0.5), sh - Inches(0.65),
                sw - Inches(1.0), Inches(0.35),
                "在 Slack 里，Viktor 不是一个机器人——是一个 7×24 的团队成员。",
                font_size=13, italic=True, color=ACCENT,
                align=PP_ALIGN.CENTER)


def build_slide_B(slide, sw, sh):
    set_slide_bg(slide, BG)

    # eyebrow
    add_textbox(slide, Inches(0.5), Inches(0.4),
                Inches(4.0), Inches(0.35),
                "07 / 延展思考",
                font_size=14, bold=True, color=ACCENT)

    # title
    add_textbox(slide, Inches(0.5), Inches(0.75),
                Inches(9.0), Inches(0.6),
                "AI 时代软件正在发生的 3 个位移",
                font_size=32, bold=True, color=WHITE)

    # subtitle
    add_textbox(slide, Inches(0.5), Inches(1.4),
                Inches(9.0), Inches(0.35),
                "Why Viktor is a glimpse of what's coming — Agentic Software",
                font_size=13, italic=True, color=SUB)

    # 3 horizontal cards
    margin_l = Inches(0.45)
    margin_r = Inches(0.45)
    top_grid = Inches(1.95)
    gap = Inches(0.25)
    grid_w = sw - margin_l - margin_r
    card_w = (grid_w - 2 * gap) / 3
    card_h = Inches(2.95)

    cards_data = [
        ("①  从「工具」到「同事」", [
            ("旧", "“我用 AI 做 X”（我主动拆任务）", SUB),
            ("新", "“让 Viktor 去做 X”（AI 自主规划）", WHITE),
            ("位移点", "主动权从人 → 部分让给 AI", ACCENT2),
        ]),
        ("②  Seat License → Task Credit", [
            ("旧", "$15/人/月（按座位）", SUB),
            ("新", "$2.5 / 1000 积分（按任务）", WHITE),
            ("位移点", "定价和 AI 执行成本直接绑定", ACCENT2),
        ]),
        ("③  人类流程 → 协作礼仪", [
            ("旧", "工作流 = 人 → 工具 → 人", SUB),
            ("新", "工作流 = 人 ↔ AI 同事 ↔ 人", WHITE),
            ("位移点", "审核边界 / 失误问责 / 响应 SLA 需要重新定义", ACCENT2),
        ]),
    ]
    for i, (title, rows) in enumerate(cards_data):
        left = margin_l + i * (card_w + gap)
        c = add_card(slide, left, top_grid, card_w, card_h)
        fill_card_compare(c, title, rows)

    # bottom note
    add_textbox(slide, Inches(0.5), sh - Inches(0.65),
                sw - Inches(1.0), Inches(0.35),
                "这不是未来，是正在发生。Viktor 只是这波浪潮里一个样本。",
                font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)


def move_slide_before(prs, slide_obj, target_index):
    """Move slide_obj in sldIdLst to be before position target_index (0-based)."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    # find current element by matching rId to slide part
    # python-pptx slide has part; we look up sldId whose rId references slide.part
    target_rid = None
    for rel in prs.part.rels.values():
        if rel.target_part is slide_obj.part:
            target_rid = rel.rId
            break
    assert target_rid is not None, "slide rId not found"
    # locate element
    current_el = None
    for el in slides_list:
        if el.get(qn('r:id')) == target_rid:
            current_el = el
            break
    assert current_el is not None
    xml_slides.remove(current_el)
    # reinsert at target_index
    # xml_slides is a list-like; use insert by XML index
    ref_el = list(xml_slides)[target_index]
    ref_el.addprevious(current_el)


def main():
    prs = Presentation(PATH)
    sw, sh = prs.slide_width, prs.slide_height

    # Use blank layout (layout 6 in default; pick the one with fewest placeholders)
    layouts = prs.slide_layouts
    blank_layout = None
    for lo in layouts:
        if len(lo.placeholders) == 0:
            blank_layout = lo
            break
    if blank_layout is None:
        # fallback — use last layout
        blank_layout = layouts[-1]

    # Add slide A at end
    slide_a = prs.slides.add_slide(blank_layout)
    build_slide_A(slide_a, sw, sh)

    # Add slide B at end
    slide_b = prs.slides.add_slide(blank_layout)
    build_slide_B(slide_b, sw, sh)

    # Now move them to before original slide 43 (index 42)
    # After adding, sldIdLst order is [0..42 original, 43=A, 44=B]
    # We want final order: [0..41 original, A, B, 42 original(the thank you page)]
    # So move slide_a to position 42, then slide_b to position 43.
    move_slide_before(prs, slide_a, 42)
    move_slide_before(prs, slide_b, 43)

    prs.save(PATH)
    print("Saved:", PATH)

    # Verify
    prs2 = Presentation(PATH)
    print("Total slides:", len(prs2.slides))
    for idx in [42, 43, 44]:
        s = prs2.slides[idx]
        texts = []
        for shp in s.shapes:
            if shp.has_text_frame:
                for para in shp.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text.strip())
        print(f"slide {idx+1} first texts:", texts[:4])


if __name__ == "__main__":
    main()
