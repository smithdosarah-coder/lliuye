# -*- coding: utf-8 -*-
"""V3: properly scaled to 10x5.625 template."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r'C:/Users/Mr.S/xwechat_files/wxid_r0ty34wac6bt22_7a6c/msg/file/2026-04/Viktor-分享pptx.pptx'
DST = r'C:/Users/Mr.S/Desktop/Viktor_v3.pptx'

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height   # 10 x 5.625 in
# Safe content area margins
MARGIN = Inches(0.35)
INNER_W = SW - 2 * MARGIN
INNER_H = SH - 2 * MARGIN

NAVY   = RGBColor(0x0A, 0x0E, 0x1A)
CARD   = RGBColor(0x14, 0x1B, 0x2D)
CARD2  = RGBColor(0x1C, 0x24, 0x3B)
BLUE   = RGBColor(0x6C, 0x8A, 0xFF)
BLUE_D = RGBColor(0x3A, 0x52, 0xB8)
ORANGE = RGBColor(0xFF, 0xB6, 0x4E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xAE, 0xB3, 0xC2)
DIM    = RGBColor(0x5A, 0x62, 0x75)

_layout = prs.slide_layouts[len(prs.slide_layouts)-1]

def blank_slide():
    s = prs.slides.add_slide(_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size=11, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ('margin_top','margin_bottom','margin_left','margin_right'):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

# ====================================================================
# Slide A — 协作机制 (左色块 + 右 feature 栈)
# Template 10 × 5.625
# ====================================================================
sA = blank_slide()

left_w = int(SW * 0.38)   # 3.8 in
left_bg = sA.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, left_w, SH)
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = CARD
left_bg.line.fill.background()

stripe = sA.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(40000), SH)
stripe.fill.solid(); stripe.fill.fore_color.rgb = BLUE
stripe.line.fill.background()

# Left pane
txt(sA, Inches(0.3), Inches(0.3), Inches(3.3), Inches(0.3),
    '06 / 协作机制', size=10, color=BLUE, bold=True)
txt(sA, Inches(0.3), Inches(1.15), Inches(3.4), Inches(1.4),
    '"AI 同事"', size=48, color=WHITE, bold=True)
txt(sA, Inches(0.3), Inches(2.35), Inches(3.4), Inches(0.6),
    '不是修辞。', size=26, color=BLUE, bold=True)
txt(sA, Inches(0.3), Inches(5.0), Inches(3.4), Inches(0.3),
    '— 它在 Slack 里真的像一个 team member',
    size=9, color=MUTED, italic=True)

# Right pane
RX = Inches(4.1)
RW = Inches(5.55)
txt(sA, RX, Inches(0.35), RW, Inches(0.45),
    '它和你用过的所有 AI 的根本区别', size=20, color=WHITE, bold=True)
txt(sA, RX, Inches(0.85), RW, Inches(0.3),
    '都在这四个地方', size=10, color=MUTED)

features = [
    ('01', '在频道里干活',         '不是私聊窗口——全组都能看到它在做什么、做到哪一步'),
    ('02', '任务对话就是知识库',   '不用再问"上次谁负责"；Slack 搜一下全过程都在'),
    ('03', '随时插手不用等',       '任何人 @ 它就能改方向、接手、喊停——像对同事下指令'),
    ('04', '会主动找你',           '周一早上："上周的任务你还要我盯吗？"'),
]
row_y0 = Inches(1.35)
row_h  = Inches(0.93)
for i, (num, title_t, desc) in enumerate(features):
    y = row_y0 + i * row_h
    txt(sA, RX, y + Inches(0.06), Inches(0.55), Inches(0.5),
        num, size=20, color=ORANGE, bold=True)
    txt(sA, RX + Inches(0.62), y + Inches(0.03), Inches(4.85), Inches(0.35),
        title_t, size=15, color=WHITE, bold=True)
    txt(sA, RX + Inches(0.62), y + Inches(0.42), Inches(4.85), Inches(0.45),
        desc, size=10, color=MUTED)
    if i < 3:
        ln = sA.shapes.add_connector(
            1, RX, y + row_h - Inches(0.05),
            RX + RW - Inches(0.1), y + row_h - Inches(0.05))
        ln.line.color.rgb = DIM
        ln.line.width = Emu(3000)

# ====================================================================
# Slide B — 一次真实任务 (时间线)
# ====================================================================
sB = blank_slide()

txt(sB, MARGIN, Inches(0.3), Inches(4), Inches(0.3),
    '07 / 实操演示', size=10, color=BLUE, bold=True)
txt(sB, MARGIN, Inches(0.62), Inches(9.3), Inches(0.55),
    '一次真实任务：从"一句话"到交付', size=24, color=WHITE, bold=True)

# quote box
pbox = sB.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           MARGIN, Inches(1.3), Inches(9.3), Inches(0.5))
pbox.fill.solid(); pbox.fill.fore_color.rgb = CARD
pbox.line.color.rgb = BLUE; pbox.line.width = Pt(0.75)
pbox.adjustments[0] = 0.3
txt(sB, MARGIN + Inches(0.2), Inches(1.4), Inches(9.0), Inches(0.35),
    '"帮我做一份近 12 个月国内 AI 编程工具市场调研，Slack 发摘要，全文存 Notion。"',
    size=11, color=WHITE, italic=True, anchor=MSO_ANCHOR.MIDDLE)

steps = [
    ('1', '接到指令',  'Slack 私聊或频道',                                    '~ 15 秒'),
    ('2', '自主规划',  '拆成 7 个子任务，列清单，反问是否开始',               '~ 30 秒'),
    ('3', '并发调工具', 'Search + Notion + 表格 + Doc，过程在频道直播',       '~ 12 分钟'),
    ('4', '交付成果',  'Slack 贴摘要链接，Notion 建好结构化页',              '~ 5 秒'),
]
# Timeline area: y=2.2 to 4.8, 2.6" height
# 4 columns within safe 9.3" → each 2.325"
timeline_y = Inches(2.35)
col_w = Inches(2.325)
circle_d = Inches(0.65)
for i, (num, title_t, desc, dur) in enumerate(steps):
    col_cx = MARGIN + col_w * i + col_w / 2
    # circle
    circ = sB.shapes.add_shape(MSO_SHAPE.OVAL,
                               col_cx - circle_d/2, timeline_y, circle_d, circle_d)
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    txt(sB, col_cx - circle_d/2, timeline_y, circle_d, circle_d,
        num, size=22, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    txt(sB, col_cx - Inches(1.05), timeline_y + Inches(0.78),
        Inches(2.1), Inches(0.35),
        title_t, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # desc
    txt(sB, col_cx - Inches(1.1), timeline_y + Inches(1.18),
        Inches(2.2), Inches(0.7),
        desc, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    # duration
    txt(sB, col_cx - Inches(0.8), timeline_y + Inches(1.95),
        Inches(1.6), Inches(0.25),
        dur, size=9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    # arrow
    if i < 3:
        ar_w = Inches(1.3)
        ar = sB.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            col_cx + circle_d/2 + Inches(0.08),
            timeline_y + circle_d/2 - Inches(0.08),
            ar_w, Inches(0.16))
        ar.fill.solid(); ar.fill.fore_color.rgb = BLUE_D
        ar.line.fill.background()

# bottom callout
bar = sB.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         MARGIN, Inches(4.85), Inches(9.3), Inches(0.55))
bar.fill.solid(); bar.fill.fore_color.rgb = CARD2
bar.line.color.rgb = ORANGE; bar.line.width = Pt(1)
bar.adjustments[0] = 0.3
txt(sB, MARGIN + Inches(0.2), Inches(4.98), Inches(9.0), Inches(0.3),
    '总耗时 ≈ 13 分钟。自己做？2–3 小时起步。',
    size=13, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# ====================================================================
# Slide C — 3 个位移 (左 hero "3" + 右 3 lanes)
# ====================================================================
sC = blank_slide()

txt(sC, MARGIN, Inches(0.3), Inches(4), Inches(0.3),
    '08 / 延展思考', size=10, color=BLUE, bold=True)

# hero "3"
txt(sC, Inches(0.15), Inches(0.75), Inches(3.7), Inches(3.4),
    '3', size=180, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.15), Inches(4.15), Inches(3.7), Inches(0.45),
    '个根本位移', size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.15), Inches(4.58), Inches(3.7), Inches(0.35),
    '正在重写所有 AI 原生产品', size=10, color=MUTED, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.15), Inches(4.98), Inches(3.7), Inches(0.3),
    'Viktor 是其中一个样本', size=9, color=BLUE, italic=True, align=PP_ALIGN.CENTER)

# right lanes
RX2 = Inches(4.15)
RW2 = Inches(5.5)
txt(sC, RX2, Inches(0.62), RW2, Inches(0.5),
    'Agentic Software 正在发生的三件事', size=18, color=WHITE, bold=True)

shifts = [
    ('工具',          '同事',        '主动权从"我拆任务"→"委派任务"'),
    ('Seat License',  'Task Credit', '定价从"按座位"→"按 AI 执行成本"'),
    ('人类流程',      '协作礼仪',    '谁审 AI 输出？失误谁负责？SLA 多少？'),
]
ly0 = Inches(1.3)
lh  = Inches(1.1)
for i, (old, new, note) in enumerate(shifts):
    y = ly0 + i * lh
    card = sC.shapes.add_shape(MSO_SHAPE.RECTANGLE, RX2, y, RW2, Inches(0.95))
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = DIM; card.line.width = Pt(0.5)
    # old
    txt(sC, RX2 + Inches(0.15), y + Inches(0.1), Inches(1.8), Inches(0.4),
        old, size=17, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrow
    ar = sC.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
        RX2 + Inches(2.05), y + Inches(0.22),
        Inches(0.45), Inches(0.17))
    ar.fill.solid(); ar.fill.fore_color.rgb = BLUE
    ar.line.fill.background()
    # new
    txt(sC, RX2 + Inches(2.6), y + Inches(0.1), Inches(2.75), Inches(0.4),
        new, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # note
    txt(sC, RX2 + Inches(0.15), y + Inches(0.55), Inches(5.2), Inches(0.35),
        note, size=9, color=ORANGE)

# ====================================================================
# Slide D — 金句收尾
# ====================================================================
sD = blank_slide()

txt(sD, MARGIN, Inches(0.3), Inches(4), Inches(0.3),
    '09 / 留个话', size=10, color=BLUE, bold=True)

txt(sD, Inches(0.3), Inches(1.8), Inches(9.4), Inches(0.9),
    'AI 时代的软件，不是被 AI 加强的软件。',
    size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sD, Inches(0.3), Inches(2.55), Inches(9.4), Inches(0.9),
    '是为 AI 协作重写的软件。',
    size=26, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# divider short line
ln = sD.shapes.add_connector(1, Inches(4.55), Inches(3.75),
                             Inches(5.45), Inches(3.75))
ln.line.color.rgb = ORANGE; ln.line.width = Pt(1.5)

txt(sD, Inches(0.3), Inches(3.95), Inches(9.4), Inches(0.5),
    'Viktor 是这波浪潮里一个样本，不是终点。',
    size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

txt(sD, Inches(0.3), Inches(4.95), Inches(9.4), Inches(0.35),
    '10,000 免费积分试试看 · 选一件你不想再重复做的事',
    size=10, color=ORANGE, align=PP_ALIGN.CENTER)

# ====================================================================
# Move 感谢观看 to end
# ====================================================================
sldIdLst = prs.slides._sldIdLst
sld_ids = list(sldIdLst)
thank = sld_ids[42]
sldIdLst.remove(thank)
sldIdLst.append(thank)

prs.save(DST)

prs2 = Presentation(DST)
print(f'Total slides: {len(prs2.slides)}  Size: {prs2.slide_width/914400} x {prs2.slide_height/914400} in')
for i in range(42, len(prs2.slides)):
    s = prs2.slides[i]
    titles = []
    for sh in s.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t and len(t) < 80:
                titles.append(t.replace('\n', ' '))
    print(f'Slide {i+1}: {titles[:3]}')
