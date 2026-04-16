# -*- coding: utf-8 -*-
"""Rebuild Viktor PPTX enhanced slides. 4 distinct-layout slides, no emoji, varied composition."""
import shutil, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r'C:/Users/Mr.S/xwechat_files/wxid_r0ty34wac6bt22_7a6c/msg/file/2026-04/Viktor-分享pptx.pptx'
DST = r'C:/Users/Mr.S/Desktop/Viktor_v2.pptx'

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height

NAVY   = RGBColor(0x0A, 0x0E, 0x1A)
CARD   = RGBColor(0x14, 0x1B, 0x2D)
CARD2  = RGBColor(0x1C, 0x24, 0x3B)
BLUE   = RGBColor(0x6C, 0x8A, 0xFF)
BLUE_D = RGBColor(0x3A, 0x52, 0xB8)
ORANGE = RGBColor(0xFF, 0xB6, 0x4E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0xAE, 0xB3, 0xC2)
DIM    = RGBColor(0x5A, 0x62, 0x75)

_layout = prs.slide_layouts[len(prs.slide_layouts)-1]
def blank_slide():
    s = prs.slides.add_slide(_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    return s

def txt(slide, x, y, w, h, text, size=14, color=WHITE, bold=False, italic=False,
        align=PP_ALIGN.LEFT, font='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ('margin_top','margin_bottom','margin_left','margin_right'):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

# ====================================================================
# Slide A — 协作机制 (左大字 + 右 4-row feature stack with dividers)
# ====================================================================
sA = blank_slide()

left_w = int(SW * 0.38)
left_bg = sA.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, left_w, SH)
left_bg.fill.solid(); left_bg.fill.fore_color.rgb = CARD
left_bg.line.fill.background()

# accent vertical stripe at left edge
stripe = sA.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(60000), SH)
stripe.fill.solid(); stripe.fill.fore_color.rgb = BLUE
stripe.line.fill.background()

txt(sA, Inches(0.7), Inches(0.6), Inches(3), Inches(0.3), '06 / 协作机制', size=12, color=BLUE, bold=True)
txt(sA, Inches(0.7), Inches(1.3), Inches(4), Inches(2),
    '"AI 同事"', size=72, color=WHITE, bold=True)
txt(sA, Inches(0.7), Inches(2.9), Inches(4), Inches(0.9),
    '不是修辞。', size=40, color=BLUE, bold=True)
txt(sA, Inches(0.7), Inches(6.7), Inches(4), Inches(0.4),
    '— 它在 Slack 里真的像一个 team member', size=12, color=MUTED, italic=True)

# right title
txt(sA, Inches(5.35), Inches(0.65), Inches(7.5), Inches(0.6),
    '它和你用过的所有 AI 的根本区别', size=26, color=WHITE, bold=True)
txt(sA, Inches(5.35), Inches(1.25), Inches(7.5), Inches(0.4),
    '都在这四个地方', size=14, color=MUTED)

features = [
    ('01', '在频道里干活', '不是私聊窗口——全组都能看到它在做什么、做到哪一步'),
    ('02', '任务对话就是知识库', '不用再问"上次谁负责"；Slack 搜索一下，全过程都在'),
    ('03', '随时插手不用等', '任何人 @ 它就能改方向、接手、喊停——像对同事下指令'),
    ('04', '会主动找你', '周一早上："上周的任务你还要我盯吗？"'),
]
row_y = Inches(2.0); row_h = Inches(1.05)
for i, (num, title_t, desc) in enumerate(features):
    y = row_y + i * row_h
    txt(sA, Inches(5.35), y + Inches(0.10), Inches(0.65), Inches(0.65),
        num, size=28, color=ORANGE, bold=True)
    txt(sA, Inches(6.05), y + Inches(0.05), Inches(6.8), Inches(0.4),
        title_t, size=20, color=WHITE, bold=True)
    txt(sA, Inches(6.05), y + Inches(0.55), Inches(6.8), Inches(0.45),
        desc, size=13, color=MUTED)
    if i < 3:
        ln = sA.shapes.add_connector(1, Inches(5.35), y + Inches(1.0),
                                     Inches(12.85), y + Inches(1.0))
        ln.line.color.rgb = DIM
        ln.line.width = Emu(4000)

# ====================================================================
# Slide B — 一次真实任务 (horizontal timeline)
# ====================================================================
sB = blank_slide()

txt(sB, Inches(0.7), Inches(0.55), Inches(4), Inches(0.3), '07 / 实操演示', size=12, color=BLUE, bold=True)
txt(sB, Inches(0.7), Inches(0.95), Inches(12), Inches(0.8),
    '一次真实任务：从"一句话"到交付', size=34, color=WHITE, bold=True)
# quote of user prompt
pbox = sB.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(0.7), Inches(1.95), Inches(12.0), Inches(0.7))
pbox.fill.solid(); pbox.fill.fore_color.rgb = CARD
pbox.line.color.rgb = BLUE; pbox.line.width = Pt(0.75)
pbox.adjustments[0] = 0.2
txt(sB, Inches(1.0), Inches(2.10), Inches(11.5), Inches(0.45),
    '"帮我做一份近 12 个月国内 AI 编程工具市场调研，Slack 发摘要，全文存 Notion。"',
    size=14, color=WHITE, italic=True)

steps = [
    ('1', '接到指令',  'Slack 私聊或频道', '~ 15 秒'),
    ('2', '自主规划',  '拆成 7 个子任务，列清单，反问你是否开始', '~ 30 秒'),
    ('3', '并发调工具', 'Search + Notion + 表格 + Google Doc，过程在频道直播', '~ 12 分钟'),
    ('4', '交付成果',  'Slack 贴摘要链接，Notion 有结构化全文页', '~ 5 秒'),
]
by = Inches(3.4)
for i, (num, title_t, desc, dur) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    circ = sB.shapes.add_shape(MSO_SHAPE.OVAL, x, by, Inches(0.85), Inches(0.85))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
    circ.line.fill.background()
    txt(sB, x, by, Inches(0.85), Inches(0.85), num,
        size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(sB, x - Inches(1.05), by + Inches(1.05), Inches(2.95), Inches(0.4),
        title_t, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sB, x - Inches(1.1), by + Inches(1.55), Inches(3.05), Inches(1.0),
        desc, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    txt(sB, x - Inches(0.6), by + Inches(2.45), Inches(2.05), Inches(0.35),
        dur, size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        ar = sB.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 x + Inches(1.05), by + Inches(0.33), Inches(2), Inches(0.22))
        ar.fill.solid(); ar.fill.fore_color.rgb = BLUE_D
        ar.line.fill.background()

# bottom callout
bar = sB.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.7))
bar.fill.solid(); bar.fill.fore_color.rgb = CARD2
bar.line.color.rgb = ORANGE; bar.line.width = Pt(1)
bar.adjustments[0] = 0.3
txt(sB, Inches(1.0), Inches(6.55), Inches(11.5), Inches(0.4),
    '总耗时 ≈ 13 分钟。自己做？2–3 小时起步。',
    size=16, color=WHITE, bold=True)

# ====================================================================
# Slide C — 3 个位移 (左大字"3" + 右 3-lane arrows)
# ====================================================================
sC = blank_slide()

txt(sC, Inches(0.7), Inches(0.55), Inches(4), Inches(0.3),
    '08 / 延展思考', size=12, color=BLUE, bold=True)

# hero number "3"
txt(sC, Inches(0.3), Inches(0.95), Inches(5.2), Inches(4.5),
    '3', size=300, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.3), Inches(5.5), Inches(5.2), Inches(0.6),
    '个根本位移', size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.3), Inches(6.05), Inches(5.2), Inches(0.5),
    '正在重写所有 AI 原生产品', size=13, color=MUTED, align=PP_ALIGN.CENTER)
txt(sC, Inches(0.3), Inches(6.7), Inches(5.2), Inches(0.4),
    'Viktor 是其中一个样本', size=12, color=BLUE, italic=True, align=PP_ALIGN.CENTER)

# right — 3 lanes
txt(sC, Inches(6.0), Inches(0.95), Inches(7.2), Inches(0.6),
    'Agentic Software 正在发生的三件事', size=25, color=WHITE, bold=True)

shifts = [
    ('工具',          '同事',        '主动权从"我拆任务"→"委派任务"'),
    ('Seat License',  'Task Credit', '定价从"按座位"→"按 AI 执行成本"'),
    ('人类流程',      '协作礼仪',    '谁审 AI 输出？失误谁负责？响应 SLA 多少？'),
]
ly = Inches(2.0); lh = Inches(1.45)
for i, (old, new, note) in enumerate(shifts):
    y = ly + i * lh
    card = sC.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), y, Inches(7.2), Inches(1.25))
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = DIM; card.line.width = Pt(0.5)
    txt(sC, Inches(6.2), y + Inches(0.22), Inches(2.3), Inches(0.55),
        old, size=22, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ar = sC.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                             Inches(8.6), y + Inches(0.38), Inches(0.6), Inches(0.22))
    ar.fill.solid(); ar.fill.fore_color.rgb = BLUE
    ar.line.fill.background()
    txt(sC, Inches(9.3), y + Inches(0.22), Inches(3.5), Inches(0.55),
        new, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(sC, Inches(6.2), y + Inches(0.82), Inches(6.8), Inches(0.4),
        note, size=11, color=ORANGE)

# ====================================================================
# Slide D — 金句收尾 (centered typography hero)
# ====================================================================
sD = blank_slide()

txt(sD, Inches(0.7), Inches(0.55), Inches(4), Inches(0.3),
    '09 / 留个话', size=12, color=BLUE, bold=True)

# big centered two-line quote
txt(sD, Inches(0.5), Inches(2.35), Inches(12.3), Inches(1.4),
    'AI 时代的软件，不是被 AI 加强的软件。',
    size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sD, Inches(0.5), Inches(3.35), Inches(12.3), Inches(1.4),
    '是为 AI 协作重写的软件。',
    size=38, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# divider short line
ln = sD.shapes.add_connector(1, Inches(6.3), Inches(5.0), Inches(7.0), Inches(5.0))
ln.line.color.rgb = ORANGE; ln.line.width = Pt(1.5)

# tagline
txt(sD, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.6),
    'Viktor 是这波浪潮里的一个样本，不是终点。',
    size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# CTA
txt(sD, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.4),
    '10,000 免费积分试试看 · 选一件你不想再重复做的事',
    size=13, color=ORANGE, align=PP_ALIGN.CENTER)

# ====================================================================
# Move 感谢观看 (original slide 43, index 42) to the very end
# ====================================================================
sldIdLst = prs.slides._sldIdLst
sld_ids = list(sldIdLst)
# Before reorder: index 42 = 感谢观看 (original), indexes 43-46 = new A/B/C/D
# We want感谢观看 to be last (index 46 after reorder)
thank = sld_ids[42]
sldIdLst.remove(thank)
sldIdLst.append(thank)

prs.save(DST)

# Verify
prs2 = Presentation(DST)
print(f'Total slides: {len(prs2.slides)}')
for i, s in enumerate(prs2.slides, 1):
    if i < 42: continue
    titles = []
    for shape in s.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 80:
                titles.append(t.replace('\n', ' '))
    print(f'Slide {i}: {titles[:3]}')
