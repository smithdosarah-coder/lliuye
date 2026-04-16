# -*- coding: utf-8 -*-
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r"C:/Users/Mr.S/xwechat_files/wxid_r0ty34wac6bt22_7a6c/msg/file/2026-04/Viktor-分享pptx.pptx"
p = Presentation(path)

def walk_shapes(shapes, indent=0):
    """Recursively extract text from all shapes including groups."""
    for shape in shapes:
        prefix = "  " * indent
        try:
            stype = shape.shape_type
        except:
            stype = "?"
        name = getattr(shape, "name", "?")

        # Group -> recurse
        if stype == MSO_SHAPE_TYPE.GROUP:
            print(f"{prefix}[GROUP:{name}]")
            walk_shapes(shape.shapes, indent+1)
            continue

        # Text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs)
                if text.strip():
                    print(f"{prefix}{text}")

        # Table
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip().replace("\n", " / ") for cell in row.cells]
                print(f"{prefix}TABLE: {' | '.join(cells)}")

        # Picture (no text but note presence)
        if stype == MSO_SHAPE_TYPE.PICTURE:
            print(f"{prefix}[PICTURE:{name}]")

        # SmartArt
        if stype == MSO_SHAPE_TYPE.DIAGRAM or "SmartArt" in name:
            print(f"{prefix}[SMARTART:{name}]")
            # try to get text from element tree
            try:
                import re
                xml = shape._element.xml
                texts = re.findall(r"<a:t>([^<]+)</a:t>", xml)
                for t in texts:
                    print(f"{prefix}  SA: {t}")
            except Exception as e:
                print(f"{prefix}  [SA extract failed: {e}]")

        # Chart
        try:
            if shape.has_chart:
                print(f"{prefix}[CHART:{name}]")
        except:
            pass

for i, slide in enumerate(p.slides, 1):
    print(f"\n========== Slide {i} ==========")
    walk_shapes(slide.shapes)
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"[NOTES]: {notes}")
